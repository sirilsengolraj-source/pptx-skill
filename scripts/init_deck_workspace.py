#!/usr/bin/env python3
"""Create a persistent workspace for iterative deck authoring."""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation

from design_tokens import PRESETS


def _slugify(value: str) -> str:
    text = re.sub(r"[^A-Za-z0-9]+", "-", value.strip().lower()).strip("-")
    return text or "deck-workspace"


def _copy_json(src: Path) -> dict[str, Any]:
    return json.loads(src.read_text(encoding="utf-8"))


def _shape_text(shape: Any) -> str:
    if not hasattr(shape, "text_frame"):
        return ""
    chunks: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        text = paragraph.text.strip()
        if text:
            chunks.append(text)
    return "\n".join(chunks).strip()


def _slide_title(slide: Any) -> str:
    try:
        title_shape = slide.shapes.title
        if title_shape is not None:
            return (title_shape.text or "").strip()
    except Exception:
        pass
    return ""


def _extract_outline(reference_pptx: Path) -> dict[str, Any]:
    prs = Presentation(str(reference_pptx))
    slides: list[dict[str, Any]] = []
    for index, slide in enumerate(prs.slides, start=1):
        title = _slide_title(slide)
        body_lines: list[str] = []
        for shape in slide.shapes:
            text = _shape_text(shape)
            if not text:
                continue
            if title and text.strip() == title:
                continue
            for line in text.splitlines():
                cleaned = line.strip()
                if cleaned:
                    body_lines.append(cleaned)
        if index == 1:
            subtitle = body_lines[0] if body_lines else ""
            slides.append({"type": "title", "title": title or "Presentation", "subtitle": subtitle})
            continue
        slide_spec: dict[str, Any] = {
            "type": "content",
            "title": title or f"Slide {index}",
        }
        if body_lines:
            slide_spec["bullets"] = body_lines[:6]
        slides.append(slide_spec)
    return {"slides": slides}


def _font_families(reference_pptx: Path) -> list[str]:
    prs = Presentation(str(reference_pptx))
    families: set[str] = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            for paragraph in shape.text_frame.paragraphs:
                paragraph_font = getattr(paragraph, "font", None)
                if paragraph_font is not None:
                    name = getattr(paragraph_font, "name", None)
                    if name:
                        families.add(name)
                for run in paragraph.runs:
                    name = getattr(run.font, "name", None)
                    if name:
                        families.add(name)
    return sorted(families)


def _reference_summary(reference_pptx: Path) -> dict[str, Any]:
    prs = Presentation(str(reference_pptx))
    return {
        "reference_pptx": str(reference_pptx),
        "slide_count": len(prs.slides),
        "slide_size_inches": {
            "width": round(prs.slide_width / 914400.0, 3),
            "height": round(prs.slide_height / 914400.0, 3),
        },
        "font_families": _font_families(reference_pptx),
    }


def _starter_outline(title: str, style_preset: str, font_pair: str | None, palette_key: str | None) -> dict[str, Any]:
    deck_style: dict[str, Any] = {
        "visual_density": "medium",
        "emoji_mode": "none",
    }
    if font_pair:
        deck_style["font_pair"] = font_pair
    if palette_key:
        deck_style["palette_key"] = palette_key
    return {
        "title": title,
        "subtitle": "Working outline",
        "deck_style": deck_style,
        "slides": [
            {
                "type": "title",
                "title": title,
                "subtitle": "Prepare notes, assets, and outline before building",
            },
            {
                "type": "content",
                "variant": "split",
                "title": "Core message",
                "subtitle": "Start from the decision or takeaway",
                "bullets": [
                    "Add the main narrative in short lines.",
                    "Keep source-backed claims explicit.",
                    "Split dense material into focused slides.",
                    "Vary layouts when the story needs rhythm.",
                ],
                "highlights": [
                    "Alignment and readability are release gates.",
                    "Titles and subtitles use dynamic vertical spacing.",
                    "Assets stay source-backed and optional.",
                    "QA blocks overlap, overflow, and sparse layouts.",
                ],
            },
        ],
    }


def _write_text(path: Path, content: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(content, encoding="utf-8")


def _workspace_readme(slug: str, title: str) -> str:
    return f"""# {title}

This workspace is the saved authoring source for the `{slug}` deck.

## Files

- `outline.json`: canonical structured slide source
- `content_plan.json`: thesis, audience, slide roles, and visual strategy
- `design_brief.json`: audience posture, cover concept, structure strategy, and grid policy
- `evidence_plan.json`: sourced claims, metrics, chart candidates, and gaps
- `style_contract.json`: stable style + layout contract for later slide additions
- `asset_plan.json`: source-backed imagery/background/chart staging plan
- `notes.md`: deck-specific data sources, decisions, and manual design notes
- `assets/`: local images, diagrams, logos, and tables used by the deck
- `build/`: generated `.pptx` output plus QA reports

## Commands

Build the deck:

```bash
python3 ../../scripts/build_workspace.py --workspace . --overwrite
```

Build and run strict QA:

```bash
python3 ../../scripts/build_workspace.py --workspace . --qa --overwrite
```

Use non-render QA when LibreOffice is unavailable:

```bash
python3 ../../scripts/build_workspace.py --workspace . --qa --skip-render --overwrite
```

Allow Wikimedia Commons fetches while staging assets:

```bash
python3 ../../scripts/build_workspace.py --workspace . --allow-network-assets --overwrite
```

## Iteration Pattern

1. Fill `content_plan.json` with thesis, audience, slide roles, and visual strategy.
2. Fill `design_brief.json` with audience posture, cover concept, and structure strategy.
3. Fill `evidence_plan.json` with sourced claims, metrics, and chart candidates.
4. Update `notes.md` with data rules and unresolved assumptions.
5. Add source-backed image/background/chart requests to `asset_plan.json`.
6. Stage local assets inside `assets/` when needed.
7. Edit `outline.json` to add, replace, or reorder slides.
8. Reference staged assets with aliases such as `asset:hero_name`, `image:crew_portrait`, or `generated:concept_visual`.
9. Re-run `build_workspace.py`.
10. Keep the source files. Do not rely on inline heredoc generation if you want to extend the deck later.
"""


def _workspace_notes(title: str, style_preset: str) -> str:
    preset = PRESETS[style_preset]
    return f"""# {title} Notes

## Purpose

- Audience:
- Decision / outcome:
- Style preset: `{style_preset}`

## Sources

- Add the datasets, URLs, or reference decks used to author this presentation.
- Record the provenance for every non-user image you stage through `asset_plan.json`.
- Promote researched claims into `evidence_plan.json` before adding them to slides.

## Research log to staging plan

Closes the gap where research produces good content but never turns into
staged visuals. Every row in this table should eventually trigger an
entry in `asset_plan.json` (wikimedia_query for a CC photo, or a staged
icon/chart).

| Fact discovered | Source | Becomes | In asset_plan as |
|---|---|---|---|
| _e.g. Chicago Pile-1, first controlled chain reaction, Dec 2 1942_ | _en.wikipedia.org/Chicago_Pile-1_ | _hero image on slide 3_ | _images[0].wikimedia_query: "Chicago Pile-1"_ |
|  |  |  |  |
|  |  |  |  |

If this table is empty at build time, ask yourself whether the deck
actually has no visual anchors or whether the research hasn't been
connected to the staging plan yet.

## Style Contract

- Slide size: 16:9 unless a reference deck says otherwise
- Title font: {preset.typography.title_max}-{preset.typography.title_min}pt range via preset
- Section font: {preset.typography.section_max}-{preset.typography.section_min}pt range via preset
- Body font: {preset.typography.body_max}-{preset.typography.body_min}pt range via preset
- Margin x: {preset.layout.margin_x}
- Gutter: {preset.layout.gutter}

## QA Notes

- Preserve alignment first.
- Keep subtitles below wrapped titles.
- Prefer local, source-backed assets in `assets/`.
- Use `asset:alias` references in `outline.json` after staging into `assets/staged/`.
- Add any deck-specific measurements here if you later match an existing deck manually.
"""


def _content_plan_stub(title: str) -> dict[str, Any]:
    return {
        "topic": title,
        "audience": "Deck author using the presentation-skill workspace scaffold.",
        "objective": "Replace the starter content with topic-specific narrative, evidence, and assets.",
        "thesis": "A reliable deck starts with a content plan, sourced evidence, staged visuals, and QA before delivery.",
        "narrative_arc": [
            {
                "act": "setup",
                "purpose": "Frame why the topic matters.",
                "slides": ["s1"],
            },
            {
                "act": "source-first setup",
                "purpose": "Show how source-first authoring keeps later edits clean without prescribing a diagram.",
                "slides": ["s2"],
            },
            {
                "act": "implication",
                "purpose": "Close with what the audience should remember or do.",
                "slides": [],
            },
        ],
        "slide_plan": [
            {
                "slide_id": "s1",
                "role": "title",
                "message": "Start from a durable workspace, not one-off inline code.",
                "variant": "title",
                "visual_strategy": "title opener with disciplined spacing",
                "evidence_needs": [],
                "asset_needs": [],
            },
            {
                "slide_id": "s2",
                "role": "setup",
                "message": "The durable source files are the contract for future edits.",
                "variant": "split",
                "visual_strategy": "two-column split with dense highlight panel",
                "evidence_needs": [],
                "asset_needs": [],
            },
        ],
        "design_notes": {
            "style_preset_reason": "Starter uses the requested preset while keeping typography and spacing conservative.",
            "rhythm_break": "Add a diagram, figure, table, or image only when the deck topic makes it useful.",
            "visual_motif": "Source-first authoring with clear stage labels when the author chooses a process visual.",
        },
    }


def _design_brief_stub(title: str, style_preset: str) -> dict[str, Any]:
    preset = PRESETS[style_preset]
    return {
        "topic": title,
        "content_maturity": "serious/work",
        "audience_posture": "coworkers/operators",
        "emotional_register": "trustworthy",
        "format_promise": (
            "A clean, editable PowerPoint deck with one dominant idea per slide, "
            "disciplined alignment, and enough visual rhythm to avoid generic card grids."
        ),
        "anti_format": [
            "repeated title plus three cards on every slide",
            "body text placed by feel instead of grid constants",
            "decorative shapes without a reading job",
            "shrinking text below readability floors to solve density",
        ],
        "canvas_and_grid": {
            "aspect": "16:9",
            "margin_x_in": 0.5,
            "footer_reserve_in": 0.32,
            "header_policy": "measured title/subtitle stack; body starts at returned contentTop",
            "column_policy": "derive columns from margin and gutter constants, not magic numbers",
        },
        "visual_system": {
            "style_preset": style_preset,
            "dominant_color": preset.palette["bg_dark"],
            "accent_primary": preset.palette["accent_primary"],
            "accent_secondary": preset.palette["accent_secondary"],
            "palette_role_map": {
                "background": "dominant or neutral field",
                "accent": "navigation, KPI emphasis, rails, and labels",
                "muted": "captions and provenance",
            },
        },
        "title_page_concept": {
            "chosen_archetype": "topic-specific opener chosen from the preset and content",
            "dominant_element": "large topic-specific title",
            "supporting_element": "short subtitle or one strong hero asset",
            "why_this_could_only_be_this_deck": "Replace with a sentence before final delivery.",
        },
        "structure_strategy": {
            "primary_scaffold": "open editorial content slides with measured headers",
            "repeated_elements": ["shared margins", "consistent source/footer treatment", "limited accent rails"],
            "allowed_variations": [
                "standard clean report slides",
                "split",
                "cards-2",
                "cards-3 with promote_card",
                "timeline only when the sequence is truly time-based",
                "table",
                "matrix",
                "optional kpi-hero only when one metric deserves isolation",
                "flow",
                "generated-image",
            ],
            "container_policy": (
                "Cards are for modular comparisons or evidence groups, not the default "
                "way to make prose look designed."
            ),
            "rhythm_break_plan": (
                "Use a rhythm break only when the content asks for it: a true "
                "hero metric, a full-bleed/source-backed image, a major section "
                "turn, or a decisive chart. Do not add a KPI hero just to break "
                "rhythm. Do not add a timeline just because a slide has steps; "
                "use report bands, a table, or a figure when those are clearer."
            ),
        },
    }


def _evidence_plan_stub(title: str) -> dict[str, Any]:
    return {
        "topic": title,
        "source_policy": "Prefer primary or source-backed facts. Do not fabricate citations.",
        "items": [],
        "chart_candidates": [],
        "open_questions": [
            "Replace the scaffold with topic-specific evidence before delivering a factual deck."
        ],
    }


def _asset_plan_stub(title: str) -> dict[str, Any]:
    """Starter plan for staged deck assets.

    The entries below are TODO placeholders (empty arrays with inline
    schema comments). They're NOT ready-to-run examples, because
    generic examples tend to ship unchanged.

    Populate the arrays with topic-specific requests, or delete the
    file if the deck doesn't need staged assets. `build_workspace.py`
    warns at build time if this file is still at its initial state and
    the deck has no icons, hero image, or charts anywhere in its
    outline.
    """
    return {
        "topic": title,
        "__readme__": (
            "Delete this __readme__ key and populate the arrays below "
            "with real image/chart/icon requests for THIS topic. "
            "See references/outline_schema.md and "
            "references/deck_workspace_mode.md for the schemas. If the "
            "deck doesn't need staged assets, delete this whole file."
        ),
        "images": [
            # Example schema - delete and replace with real entries:
            # {"name": "hero_photo", "wikimedia_query": "<topic keyword>",
            #  "allow_sharealike": true, "attribution_line": "<caption>"}
        ],
        "backgrounds": [
            # Example schema - delete and replace:
            # {"name": "section_bg", "path": "assets/staged/section_bg.png"}
        ],
        "charts": [
            # Example schema - delete and replace with real data:
            # {"name": "trend_by_year", "title": "<chart title>",
            #  "type": "line" | "bar" | "pie",
            #  "series": [{"name": "<series>", "labels": [...], "values": [...]}],
            #  "options": {"catAxisTitle": "...", "valAxisTitle": "..."}}
        ],
        "generated_images": [
            # Example schema - delete and replace with deliberate generated
            # concept art only when source-backed imagery is insufficient:
            # {"name": "concept_visual",
            #  "prompt": "A precise editorial illustration of ...",
            #  "purpose": "Optional visual anchor slide",
            #  "model": "gpt-image-2",
            #  "size": "1536x1024",
            #  "quality": "medium"}
        ],
        "icons": [
            # Example schema - delete and replace with real icon names.
            # Bare names resolve against <workspace>/assets/icons/<name>.png.
            # Use when a cards-3 / timeline / matrix / stats / cards-2
            # slide has a clear visual metaphor per card.
            # {"name": "reactor_core", "path": "assets/icons/reactor_core.png"}
        ],
    }


def _style_contract(
    *,
    title: str,
    slug: str,
    style_preset: str,
    font_pair: str | None,
    palette_key: str | None,
    reference_pptx: Path | None,
) -> dict[str, Any]:
    preset = PRESETS[style_preset]
    contract: dict[str, Any] = {
        "workspace_version": 1,
        "deck_title": title,
        "deck_slug": slug,
        "build": {
            "style_preset": style_preset,
            "font_pair": font_pair,
            "palette_key": palette_key,
            "output_pptx": f"build/{slug}.pptx",
            "qa_dir": "build/qa",
            "qa_report": "build/qa/report.json",
        },
        "layout_rules": {
            "alignment_first": True,
            "zero_overlap_required": True,
            "title_subtitle_stack_dynamic": True,
            "footer_safe_region_required": True,
            "cards_use_measured_text_fit": True,
        },
        "preset_tokens": preset.to_dict(),
    }
    if reference_pptx:
        contract["reference"] = _reference_summary(reference_pptx)
    return contract


def _args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Initialize a persistent PPTX deck workspace.")
    parser.add_argument("--workspace", required=True, help="Workspace directory to create")
    parser.add_argument("--title", required=True, help="Human-readable deck title")
    parser.add_argument("--style-preset", default="executive-clinical", choices=sorted(PRESETS))
    parser.add_argument("--font-pair", help="Optional font pair override stored in outline/style contract")
    parser.add_argument("--palette-key", help="Optional palette override stored in outline/style contract")
    parser.add_argument("--source-outline", help="Optional JSON outline to copy into workspace")
    parser.add_argument("--reference-pptx", help="Optional reference deck to summarize and extract")
    parser.add_argument("--overwrite", action="store_true", help="Replace existing workspace files")
    parser.add_argument(
        "--followup-edit",
        action="store_true",
        help=(
            "Explicitly acknowledge that --source-outline / --reference-pptx "
            "points at an existing deck and this is a followup edit to the "
            "SAME topic (not a new deck). Required when sourcing from another "
            "decks/<slug>/ directory. Prevents the 'clone-an-existing-deck-"
            "as-house-style' anti-pattern - see references/codex_guardrails.md."
        ),
    )
    return parser.parse_args()


def _is_under_decks_dir(path: Path) -> tuple[bool, str | None]:
    """Return (True, sibling_slug) if `path` lives under a decks/<slug>/ tree
    that is NOT the workspace being initialized. sibling_slug is the name
    of the source workspace, used for the error message.
    """
    resolved = path.resolve()
    for ancestor in resolved.parents:
        if ancestor.name == "decks" and ancestor.parent.name in {"presentation-skill", "pptx-skill"}:
            # Immediate child of decks/ is the source workspace slug.
            try:
                rel = resolved.relative_to(ancestor)
                return True, rel.parts[0] if rel.parts else None
            except ValueError:
                return False, None
    return False, None


def main() -> int:
    args = _args()
    workspace = Path(args.workspace).expanduser().resolve()
    source_outline = Path(args.source_outline).expanduser().resolve() if args.source_outline else None
    reference_pptx = Path(args.reference_pptx).expanduser().resolve() if args.reference_pptx else None

    if source_outline and not source_outline.exists():
        raise FileNotFoundError(f"Source outline not found: {source_outline}")
    if reference_pptx and not reference_pptx.exists():
        raise FileNotFoundError(f"Reference deck not found: {reference_pptx}")
    if workspace.exists() and any(workspace.iterdir()) and not args.overwrite:
        raise FileExistsError(f"Workspace already exists and is not empty: {workspace}")

    # Guardrail: cloning an existing deck workspace as a "house style" for a
    # new topic is a documented Codex anti-pattern (see codex_guardrails.md
    # "Eighth Trap"). If --source-outline or --reference-pptx points into
    # another decks/<slug>/ tree, require --followup-edit to acknowledge
    # that the intent is editing the same topic, not cloning style.
    for source_label, source in (
        ("--source-outline", source_outline),
        ("--reference-pptx", reference_pptx),
    ):
        if source is None:
            continue
        under_decks, source_slug = _is_under_decks_dir(source)
        if not under_decks:
            continue
        workspace_slug = _slugify(args.title)
        # If the source is in the SAME workspace we're re-initializing
        # (e.g., overwriting), let it pass. The anti-pattern is
        # cross-topic cloning.
        if source_slug == workspace_slug or source_slug == workspace.name:
            continue
        if not args.followup_edit:
            # Softened from hard ERROR to a warning: studying a past deck
            # for file shape or vocabulary is legitimate; cloning its
            # variant mix wholesale is not. Trust the author.
            print(
                f"[init_deck_workspace] WARNING: {source_label} points at "
                f"workspace {source_slug!r}. If you're just reading it for "
                "file shape, fine. If you're cloning its variant mix for a "
                "new topic, reconsider - see codex_guardrails.md on uniform-"
                "deck syndrome. Pass --followup-edit to silence this warning.",
                file=sys.stderr,
            )

    workspace.mkdir(parents=True, exist_ok=True)
    for subdir in ("assets", "assets/diagrams", "build"):
        (workspace / subdir).mkdir(parents=True, exist_ok=True)
    _write_text(workspace / "assets" / ".gitkeep", "")
    _write_text(workspace / "build" / ".gitkeep", "")

    slug = _slugify(args.title)
    if source_outline:
        outline = _copy_json(source_outline)
    elif reference_pptx:
        outline = _extract_outline(reference_pptx)
    else:
        outline = _starter_outline(args.title, args.style_preset, args.font_pair, args.palette_key)

    outline.setdefault("title", args.title)
    if args.font_pair or args.palette_key:
        deck_style = outline.setdefault("deck_style", {})
        if args.font_pair:
            deck_style.setdefault("font_pair", args.font_pair)
        if args.palette_key:
            deck_style.setdefault("palette_key", args.palette_key)

    style_contract = _style_contract(
        title=args.title,
        slug=slug,
        style_preset=args.style_preset,
        font_pair=args.font_pair,
        palette_key=args.palette_key,
        reference_pptx=reference_pptx,
    )
    workspace_manifest = {
        "workspace_version": 1,
        "deck_title": args.title,
        "deck_slug": slug,
        "style_contract": "style_contract.json",
        "content_plan": "content_plan.json",
        "design_brief": "design_brief.json",
        "evidence_plan": "evidence_plan.json",
        "outline": "outline.json",
        "asset_plan": "asset_plan.json",
        "notes": "notes.md",
        "assets_dir": "assets",
        "staged_assets_dir": "assets/staged",
        "build_dir": "build",
        "reference_pptx": str(reference_pptx) if reference_pptx else None,
    }

    _write_text(workspace / "outline.json", json.dumps(outline, indent=2, ensure_ascii=False) + "\n")
    _write_text(
        workspace / "style_contract.json",
        json.dumps(style_contract, indent=2, ensure_ascii=False) + "\n",
    )
    _write_text(
        workspace / "content_plan.json",
        json.dumps(_content_plan_stub(args.title), indent=2, ensure_ascii=False) + "\n",
    )
    _write_text(
        workspace / "design_brief.json",
        json.dumps(_design_brief_stub(args.title, args.style_preset), indent=2, ensure_ascii=False) + "\n",
    )
    _write_text(
        workspace / "evidence_plan.json",
        json.dumps(_evidence_plan_stub(args.title), indent=2, ensure_ascii=False) + "\n",
    )
    _write_text(workspace / "asset_plan.json", json.dumps(_asset_plan_stub(args.title), indent=2) + "\n")
    _write_text(workspace / "workspace.json", json.dumps(workspace_manifest, indent=2) + "\n")
    _write_text(workspace / "README.md", _workspace_readme(slug, args.title))
    _write_text(workspace / "notes.md", _workspace_notes(args.title, args.style_preset))

    print(f"Workspace created: {workspace}")
    print(f"Outline: {workspace / 'outline.json'}")
    print(f"Style contract: {workspace / 'style_contract.json'}")
    print(f"Build target: {workspace / 'build' / f'{slug}.pptx'}")
    return 0


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except Exception as exc:  # pragma: no cover - CLI error path
        print(f"Error: {exc}")
        raise SystemExit(1)
