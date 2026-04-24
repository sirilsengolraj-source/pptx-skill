#!/usr/bin/env python3
"""Run strict QA checks for PPTX quality and layout coherence."""

from __future__ import annotations

import argparse
import json
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Any

from pptx import Presentation


def _run_capture(cmd: list[str]) -> tuple[int, str]:
    result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True)
    return result.returncode, result.stdout


def _run(cmd: list[str]) -> str:
    code, output = _run_capture(cmd)
    if code != 0:
        raise RuntimeError(output.strip() or f"Command failed: {' '.join(cmd)}")
    return output


def _load_json_text(raw: str, default: Any) -> Any:
    text = raw.strip()
    if not text:
        return default
    return json.loads(text)


def _visual_summary(payload: list[dict[str, Any]]) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    warnings = [item for item in payload if item.get("severity") == "warning"]
    infos = [item for item in payload if item.get("severity") == "info"]
    return warnings, infos


def _design_summary(payload: dict[str, Any]) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    issues = payload.get("issues", [])
    warnings = [item for item in issues if item.get("severity") == "warning"]
    errors = [item for item in issues if item.get("severity") == "error"]
    return errors, warnings


def _font_families(pptx_path: Path) -> set[str]:
    families: set[str] = set()
    prs = Presentation(str(pptx_path))
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            for paragraph in shape.text_frame.paragraphs:
                paragraph_font = getattr(paragraph, "font", None)
                if paragraph_font is not None:
                    paragraph_font_name = getattr(paragraph_font, "name", None)
                    if paragraph_font_name:
                        families.add(paragraph_font_name)
                for run in paragraph.runs:
                    font_name = getattr(run.font, "name", None)
                    if font_name:
                        families.add(font_name)
    return families


def _load_json(path: Path) -> Any:
    if not path.exists():
        return {}
    return json.loads(path.read_text(encoding="utf-8"))


def _issue_summary(issues_payload: dict[str, Any]) -> tuple[int, int, int, int]:
    issue_slide_count = len(issues_payload)
    issue_shape_count = 0
    overflow_count = 0
    overlap_count = 0
    for slide_shapes in issues_payload.values():
        if not isinstance(slide_shapes, dict):
            continue
        issue_shape_count += len(slide_shapes)
        for shape_data in slide_shapes.values():
            if not isinstance(shape_data, dict):
                continue
            if "overflow" in shape_data:
                overflow_count += 1
            if "overlap" in shape_data:
                overlap_count += 1
    return issue_slide_count, issue_shape_count, overflow_count, overlap_count


def _flatten_geometry(layout_payload: dict[str, Any]) -> list[dict[str, Any]]:
    flattened: list[dict[str, Any]] = []
    for slide in layout_payload.get("slides", []):
        slide_index = slide.get("slide_index")
        for violation in slide.get("violations", []):
            item = dict(violation)
            item["slide_index"] = slide_index
            flattened.append(item)
    return flattened


def _placeholder_hits(text_path: Path) -> list[str]:
    if not text_path.exists():
        return []
    text = text_path.read_text(encoding="utf-8", errors="ignore")
    patterns = [
        r"\bxxxx\b",
        r"\blorem\b",
        r"\bipsum\b",
        r"this.*(page|slide).*layout",
    ]
    hits: list[str] = []
    for pattern in patterns:
        if re.search(pattern, text, re.IGNORECASE | re.DOTALL):
            hits.append(pattern)
    return hits


def _args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="QA gate for PPTX visual coherence.")
    parser.add_argument("--input", required=True, help="Input .pptx file")
    parser.add_argument(
        "--outdir",
        help="Output directory for QA artifacts (default: ephemeral temp dir)",
    )
    parser.add_argument(
        "--style-preset",
        default="executive-clinical",
        help="Style preset used for geometry and density thresholds",
    )
    parser.add_argument(
        "--max-font-families",
        type=int,
        default=3,
        help="Maximum allowed distinct font families (default: 3)",
    )
    parser.add_argument(
        "--max-density",
        type=float,
        help="Override max layout density threshold",
    )
    parser.add_argument(
        "--max-empty-ratio",
        type=float,
        help="Override max empty-area threshold",
    )
    parser.add_argument(
        "--max-loops",
        type=int,
        default=1,
        help="Metadata field for orchestrated loop runs (default: 1)",
    )
    parser.add_argument(
        "--outline",
        help=(
            "Optional path to the outline JSON so layout_lint can apply "
            "slide-type-aware checks (title/section empty-ratio, stats "
            "value sanity)."
        ),
    )
    parser.add_argument(
        "--strict-stats",
        action="store_true",
        help="Escalate layout_lint stats_value_non_numeric from warning to error.",
    )
    parser.add_argument(
        "--strict-geometry",
        action="store_true",
        help="Fail on error-severity geometry violations",
    )
    parser.add_argument(
        "--fail-on-geometry-warnings",
        action="store_true",
        help="In strict mode, also fail warning-severity geometry violations",
    )
    parser.add_argument(
        "--skip-render",
        action="store_true",
        help="Skip rendering slides to images",
    )
    parser.add_argument(
        "--allow-issues",
        action="store_true",
        help="Do not fail when inventory issues are found",
    )
    parser.add_argument(
        "--allow-placeholders",
        action="store_true",
        help="Do not fail when placeholder marker patterns are detected",
    )
    parser.add_argument(
        "--manual-flag",
        help=(
            "Path to a manual-review completion file. "
            "If omitted, defaults to <outdir>/manual_review_passed.flag."
        ),
    )
    parser.add_argument(
        "--skip-manual-review",
        action="store_true",
        help="Do not fail strict mode when manual-review flag is missing",
    )
    parser.add_argument(
        "--report",
        help="Path to write machine-readable QA summary JSON",
    )
    parser.add_argument(
        "--fail-on-visual-warnings",
        action="store_true",
        help="Fail when visual QA flags underfilled or sparse compositions",
    )
    parser.add_argument(
        "--fail-on-design-warnings",
        action="store_true",
        help="Fail when targeted design QA emits warning-level findings",
    )
    parser.add_argument(
        "--keep-artifacts",
        action="store_true",
        help="Keep QA artifact directory when --outdir is not provided",
    )
    return parser.parse_args()


def main() -> int:
    args = _args()
    input_path = Path(args.input).expanduser().resolve()
    if not input_path.exists():
        raise FileNotFoundError(f"Input file not found: {input_path}")

    cleanup_artifacts = False
    if args.outdir:
        outdir = Path(args.outdir).expanduser().resolve()
    else:
        outdir = Path(tempfile.mkdtemp(prefix="pptx-qa-")).resolve()
        cleanup_artifacts = not args.keep_artifacts and not args.report
    outdir.mkdir(parents=True, exist_ok=True)

    issues_json = outdir / "issues.json"
    outline_md = outdir / "outline.md"
    render_dir = outdir / "renders"
    render_dir.mkdir(parents=True, exist_ok=True)
    layout_report = outdir / "layout_lint.json"
    visual_report = outdir / "visual_qa.json"
    design_report = outdir / "design_rules.json"
    report_path = (
        Path(args.report).expanduser().resolve() if args.report else outdir / "qa_report.json"
    )
    manual_flag = (
        Path(args.manual_flag).expanduser().resolve()
        if args.manual_flag
        else outdir / "manual_review_passed.flag"
    )

    py = sys.executable
    base = Path(__file__).resolve().parent

    _run([py, str(base / "inventory.py"), str(input_path), str(issues_json), "--issues-only"])
    _run(
        [
            py,
            str(base / "extract_outline.py"),
            "--input",
            str(input_path),
            "--format",
            "markdown",
            "--output",
            str(outline_md),
        ]
    )

    lint_cmd = [
        py,
        str(base / "layout_lint.py"),
        "--input",
        str(input_path),
        "--style-preset",
        args.style_preset,
        "--output",
        str(layout_report),
    ]
    if args.max_density is not None:
        lint_cmd.extend(["--max-density", str(args.max_density)])
    if args.max_empty_ratio is not None:
        lint_cmd.extend(["--max-empty-ratio", str(args.max_empty_ratio)])
    if args.outline:
        lint_cmd.extend(["--outline", str(Path(args.outline).expanduser().resolve())])
    if args.strict_stats:
        lint_cmd.append("--strict-stats")
    _run(lint_cmd)

    if not args.skip_render:
        _run(
            [
                py,
                str(base / "render_slides.py"),
                "--input",
                str(input_path),
                "--outdir",
                str(render_dir),
                "--dpi",
                "180",
                "--format",
                "jpeg",
            ]
        )
    visual_rc, visual_out = _run_capture(
        [py, str(base / "visual_qa.py"), "--input", str(input_path), "--json"]
    )
    design_rc, _ = _run_capture(
        [
            py,
            str(base / "design_rules_qa.py"),
            "--input",
            str(input_path),
            "--report",
            str(design_report),
        ]
    )
    visual_payload = _load_json_text(visual_out, [])
    visual_report.write_text(json.dumps(visual_payload, indent=2), encoding="utf-8")
    design_payload = _load_json(design_report)

    issues_payload = _load_json(issues_json)
    layout_payload = _load_json(layout_report)
    placeholder_hits = _placeholder_hits(outline_md)
    issue_slides, issue_shapes, overflow_count, overlap_count = _issue_summary(issues_payload)
    geometry_violations = _flatten_geometry(layout_payload)
    geometry_errors = [item for item in geometry_violations if item.get("severity") == "error"]
    geometry_warnings = [
        item for item in geometry_violations if item.get("severity") == "warning"
    ]
    visual_warnings, visual_infos = _visual_summary(visual_payload)
    design_errors, design_warnings = _design_summary(design_payload)

    families = sorted(_font_families(input_path))
    too_many_fonts = len(families) > args.max_font_families
    manual_review_passed = manual_flag.exists()

    density_score_by_slide = layout_payload.get("summary", {}).get("density_score_by_slide", [])

    payload = {
        "input": str(input_path),
        "outdir": str(outdir),
        "style_preset": args.style_preset,
        "max_loops": args.max_loops,
        "overflow_count": overflow_count,
        "overlap_count": overlap_count,
        "placeholder_hits": placeholder_hits,
        "issue_slide_count": issue_slides,
        "issue_shape_count": issue_shapes,
        "geometry_violations": geometry_violations,
        "geometry_error_count": len(geometry_errors),
        "geometry_warning_count": len(geometry_warnings),
        "visual_warning_count": len(visual_warnings),
        "visual_info_count": len(visual_infos),
        "visual_report": str(visual_report),
        "visual_rc": visual_rc,
        "design_error_count": len(design_errors),
        "design_warning_count": len(design_warnings),
        "design_report": str(design_report),
        "design_rc": design_rc,
        "density_score_by_slide": density_score_by_slide,
        "font_families": families,
        "manual_review_passed": manual_review_passed,
        "strict_geometry": args.strict_geometry,
    }
    if not cleanup_artifacts or args.report:
        report_path.parent.mkdir(parents=True, exist_ok=True)
        report_path.write_text(json.dumps(payload, indent=2), encoding="utf-8")

    print(f"QA artifacts: {outdir}")
    print(f"Issue slides: {issue_slides}")
    print(f"Issue shapes: {issue_shapes}")
    print(f"Overflow count: {overflow_count}")
    print(f"Overlap count: {overlap_count}")
    print(f"Placeholder hits: {len(placeholder_hits)}")
    print(f"Geometry violations: {len(geometry_violations)}")
    print(f"Visual warnings: {len(visual_warnings)}")
    print(f"Design errors/warnings: {len(design_errors)}/{len(design_warnings)}")
    print(f"Font families ({len(families)}): {', '.join(families) if families else 'none'}")
    print(f"Manual review flag: {'present' if manual_review_passed else 'missing'} ({manual_flag})")
    if not cleanup_artifacts or args.report:
        print(f"QA report: {report_path}")
    else:
        print("QA report: ephemeral (not written, artifacts will be removed)")

    failed = False
    if issue_shapes > 0 and not args.allow_issues:
        print("FAIL: inventory detected overflow/overlap issues.")
        failed = True
    if placeholder_hits and not args.allow_placeholders:
        print("FAIL: placeholder marker patterns were detected in extracted content.")
        failed = True
    if too_many_fonts:
        print(
            f"FAIL: font coherence check exceeded threshold "
            f"({len(families)} > {args.max_font_families})."
        )
        failed = True
    if args.strict_geometry and geometry_errors:
        print("FAIL: strict geometry mode found error-level layout violations.")
        failed = True
    if args.strict_geometry and args.fail_on_geometry_warnings and geometry_warnings:
        print("FAIL: strict geometry mode found warning-level layout violations.")
        failed = True
    if design_errors:
        print("FAIL: design rules QA found error-level issues.")
        failed = True
    if args.fail_on_design_warnings and design_warnings:
        print("FAIL: design rules QA found warning-level issues.")
        failed = True
    if args.fail_on_visual_warnings and visual_warnings:
        print("FAIL: visual QA found underfilled or sparse layouts.")
        failed = True
    if args.strict_geometry and not args.skip_manual_review and not manual_review_passed:
        print("FAIL: strict geometry mode requires manual review flag.")
        failed = True

    if cleanup_artifacts:
        shutil.rmtree(outdir, ignore_errors=True)

    return 1 if failed else 0


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except Exception as exc:  # pragma: no cover - CLI error path
        print(f"Error: {exc}")
        raise SystemExit(1)
