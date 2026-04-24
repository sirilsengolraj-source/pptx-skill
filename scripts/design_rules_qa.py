#!/usr/bin/env python3
"""Targeted QA for layout-polish issues that generic geometry checks miss."""

from __future__ import annotations

import argparse
import json
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation

NS = {"c": "http://schemas.openxmlformats.org/drawingml/2006/chart"}
DEFAULT_BANNED = [
    "follow previous tool",
    "external pptx",
    "updated skill",
    "demo deck",
    "sample deck",
    "placeholder",
]


def _box(shape):
    return (
        shape.left.inches,
        shape.top.inches,
        shape.width.inches,
        shape.height.inches,
    )


def _contains(outer, inner, pad=0.02):
    ox, oy, ow, oh = outer
    ix, iy, iw, ih = inner
    return (
        ix >= ox + pad
        and iy >= oy + pad
        and ix + iw <= ox + ow - pad
        and iy + ih <= oy + oh - pad
    )


def _overlap(a, b):
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    overlap_x = min(ax + aw, bx + bw) - max(ax, bx)
    overlap_y = min(ay + ah, by + bh) - max(ay, by)
    return max(0.0, overlap_x), max(0.0, overlap_y)


def _center(box):
    x, y, w, h = box
    return (x + w / 2.0, y + h / 2.0)


def _shape_text(shape):
    if not getattr(shape, "has_text_frame", False):
        return ""
    return (shape.text or "").strip()


def _shape_kind(shape):
    try:
        if not hasattr(shape, "auto_shape_type") or shape.auto_shape_type is None:
            return ""
        return str(shape.auto_shape_type).upper()
    except Exception:
        return ""


def _has_visible_fill(shape):
    try:
        return getattr(shape.fill, "type", None) is not None
    except Exception:
        return False


def _iter_text_shapes(slide):
    for idx, shape in enumerate(slide.shapes, start=1):
        text = _shape_text(shape)
        if text:
            yield idx, shape, text


def _iter_auto_shapes(slide):
    for idx, shape in enumerate(slide.shapes, start=1):
        kind = _shape_kind(shape)
        if kind:
            yield idx, shape, kind


def check_branding(slide_idx, text_shapes, banned):
    issues = []
    for shape_id, _, text in text_shapes:
        lowered = text.lower()
        for phrase in banned:
            if phrase in lowered:
                issues.append(
                    {
                        "slide_index": slide_idx,
                        "shape_id": f"shape-{shape_id}",
                        "type": "residual_branding",
                        "severity": "error",
                        "text": text[:160],
                        "phrase": phrase,
                    }
                )
    return issues


def check_footer_overlap(slide_idx, text_shapes):
    issues = []
    bottom_band = [
        (shape_id, shape, text)
        for shape_id, shape, text in text_shapes
        if shape.top.inches >= 6.9
    ]
    for i in range(len(bottom_band)):
        for j in range(i + 1, len(bottom_band)):
            left = bottom_band[i]
            right = bottom_band[j]
            overlap_x, overlap_y = _overlap(_box(left[1]), _box(right[1]))
            if overlap_x > 0.02 and overlap_y > 0.02:
                issues.append(
                    {
                        "slide_index": slide_idx,
                        "shape_ids": [f"shape-{left[0]}", f"shape-{right[0]}"],
                        "type": "footer_text_overlap",
                        "severity": "error",
                        "delta_inches": round(min(overlap_x, overlap_y), 3),
                    }
                )
    return issues


def check_stacked_text_gaps(slide_idx, auto_shapes, text_shapes):
    issues = []
    text_boxes = [(shape_id, shape, text, _box(shape)) for shape_id, shape, text in text_shapes]
    for shape_id, shape, kind in auto_shapes:
        if "RECTANGLE" not in kind:
            continue
        box = _box(shape)
        if box[1] >= 6.7:
            continue
        if not (0.6 <= box[3] <= 4.2 and 1.0 <= box[2] <= 7.0):
            continue
        inside = [item for item in text_boxes if _contains(box, item[3])]
        if len(inside) < 2:
            continue
        inside.sort(key=lambda item: item[3][1])
        for current, nxt in zip(inside, inside[1:]):
            current_box = current[3]
            next_box = nxt[3]
            gap = next_box[1] - (current_box[1] + current_box[3])
            if gap < 0.08:
                issues.append(
                    {
                        "slide_index": slide_idx,
                        "shape_ids": [f"shape-{current[0]}", f"shape-{nxt[0]}"],
                        "container_shape_id": f"shape-{shape_id}",
                        "type": "stack_gap_too_small",
                        "severity": "warning",
                        "delta_inches": round(0.08 - gap, 3),
                    }
                )
                break
    return issues


def check_marker_centering(slide_idx, auto_shapes, text_shapes):
    issues = []
    text_candidates = [
        (shape_id, shape, text, _box(shape))
        for shape_id, shape, text in text_shapes
        if len(text.strip()) <= 3
    ]
    for shape_id, shape, kind in auto_shapes:
        if "ELLIPSE" not in kind:
            continue
        box = _box(shape)
        if not (0.2 <= box[2] <= 0.8 and abs(box[2] - box[3]) <= 0.08):
            continue
        circle_center = _center(box)
        matches = []
        for text_id, _, _, text_box in text_candidates:
            text_center = _center(text_box)
            if abs(text_center[0] - circle_center[0]) <= 0.35 and abs(text_center[1] - circle_center[1]) <= 0.35:
                matches.append((text_id, text_box, text_center))
        if not matches:
            continue
        text_id, text_box, text_center = min(
            matches,
            key=lambda item: abs(item[2][0] - circle_center[0]) + abs(item[2][1] - circle_center[1]),
        )
        dx = abs(text_center[0] - circle_center[0])
        dy = abs(text_center[1] - circle_center[1])
        if dx > 0.03 or dy > 0.03:
            issues.append(
                {
                    "slide_index": slide_idx,
                    "shape_ids": [f"shape-{shape_id}", f"shape-{text_id}"],
                    "type": "marker_label_off_center",
                    "severity": "error",
                    "delta_inches": round(max(dx, dy), 3),
                }
            )
    return issues


def check_chart_headroom(pptx_path: Path):
    issues = []
    with zipfile.ZipFile(pptx_path, "r") as archive:
        for name in archive.namelist():
            if not name.startswith("ppt/charts/chart") or not name.endswith(".xml"):
                continue
            root = ET.fromstring(archive.read(name))
            show_val = any(
                node.attrib.get("val") in {"1", "true", "True"}
                for node in root.findall(".//c:dLbls/c:showVal", NS)
            )
            if not show_val:
                continue
            axis_max_values = [
                float(node.attrib.get("val"))
                for node in root.findall(".//c:valAx/c:scaling/c:max", NS)
                if node.attrib.get("val")
            ]
            if not axis_max_values:
                continue
            axis_max = max(axis_max_values)
            point_values = [
                float(node.text)
                for node in root.findall(".//c:ser//c:val//c:v", NS)
                if node.text
            ]
            if point_values and max(point_values) >= axis_max:
                issues.append(
                    {
                        "chart_part": name,
                        "type": "chart_value_label_headroom_risk",
                        "severity": "warning",
                        "axis_max": axis_max,
                        "max_value": max(point_values),
                    }
                )
    return issues


def main() -> int:
    parser = argparse.ArgumentParser(description="Targeted design QA")
    parser.add_argument("--input", required=True, help="Input PPTX")
    parser.add_argument("--report", help="Optional JSON report path")
    parser.add_argument(
        "--banned-phrase",
        action="append",
        default=[],
        help="Additional banned phrase to flag",
    )
    args = parser.parse_args()

    pptx_path = Path(args.input).expanduser().resolve()
    prs = Presentation(str(pptx_path))

    banned = [item.lower() for item in (DEFAULT_BANNED + args.banned_phrase)]
    issues = []
    slide_summaries = []

    for slide_idx, slide in enumerate(prs.slides):
        text_shapes = list(_iter_text_shapes(slide))
        auto_shapes = list(_iter_auto_shapes(slide))
        slide_issues = []
        slide_issues.extend(check_branding(slide_idx, text_shapes, banned))
        slide_issues.extend(check_footer_overlap(slide_idx, text_shapes))
        slide_issues.extend(check_stacked_text_gaps(slide_idx, auto_shapes, text_shapes))
        slide_issues.extend(check_marker_centering(slide_idx, auto_shapes, text_shapes))
        issues.extend(slide_issues)
        slide_summaries.append({"slide_index": slide_idx, "issue_count": len(slide_issues)})

    chart_issues = check_chart_headroom(pptx_path)
    issues.extend(chart_issues)

    payload = {
        "input": str(pptx_path),
        "issue_count": len(issues),
        "error_count": sum(1 for item in issues if item.get("severity") == "error"),
        "warning_count": sum(1 for item in issues if item.get("severity") == "warning"),
        "slides": slide_summaries,
        "issues": issues,
        "passed": not issues,
    }

    if args.report:
        report_path = Path(args.report).expanduser().resolve()
        report_path.parent.mkdir(parents=True, exist_ok=True)
        report_path.write_text(json.dumps(payload, indent=2), encoding="utf-8")

    print(
        f"Design rules QA: {pptx_path}\n"
        f"  {payload['issue_count']} issue(s) | "
        f"errors={payload['error_count']} warnings={payload['warning_count']}"
    )
    for issue in issues:
        location = (
            f"slide {issue.get('slide_index', 0) + 1}"
            if "slide_index" in issue
            else issue.get("chart_part", "chart")
        )
        print(f"  - {location}: {issue.get('type')}")

    return 0 if payload["passed"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
